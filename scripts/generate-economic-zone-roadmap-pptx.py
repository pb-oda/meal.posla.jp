#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
POSLA 経済圏ロードマップ プレゼンテーション生成スクリプト
対象: 共同経営者・社内チーム + 投資家・銀行
構成: 3段階 (今 / 中間 / 将来) のロードマップを軸に、
      戦略意図 + KPI + ユニットエコノミクス + リスクまで網羅
出力: docs/posla-economic-zone-roadmap.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# デザイン定数
# ============================================================
NAVY = RGBColor(0x1A, 0x23, 0x7E)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x15, 0x65, 0xC0)

FONT_JP = 'Hiragino Sans'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOTAL_SLIDES = 20

# ============================================================
# ヘルパー関数
# ============================================================

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_footer(slide, page_num):
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.38), Inches(8), Inches(0.35),
             'POSLA 経済圏ロードマップ  |  Plus Belief Inc.  |  社外秘',
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.38), Inches(1.1), Inches(0.35),
             '%d / %d' % (page_num, TOTAL_SLIDES), size=10, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_title_bar(slide, title, page_num, subtitle=None):
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.0), NAVY)
    add_rect(slide, Emu(0), Inches(1.0), SLIDE_W, Inches(0.08), ORANGE)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             title, size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, SLIDE_W - Inches(4.5), Inches(0.32), Inches(4.0), Inches(0.45),
                 subtitle, size=14, color=ORANGE, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, page_num)


def add_bullet_list(slide, left, top, width, height, items, size=16, line_spacing=1.4,
                    mark_color=ORANGE, text_color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        run_mark = p.add_run()
        run_mark.text = '■  '
        run_mark.font.name = FONT_JP
        run_mark.font.size = Pt(size)
        run_mark.font.color.rgb = mark_color
        run_mark.font.bold = True

        run = p.add_run()
        run.text = item
        run.font.name = FONT_JP
        run.font.size = Pt(size)
        run.font.color.rgb = text_color
    return box


def add_phase_card(slide, left, top, width, height, badge, title, period,
                   tenants, mrr, body_lines, accent_color=NAVY):
    """3段階ロードマップ用のフェーズカード"""
    add_rect(slide, left + Emu(50000), top + Emu(50000), width, height,
             RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, left, top, width, height, WHITE, line_color=accent_color)
    add_rect(slide, left, top, width, Inches(0.85), accent_color)
    # バッジ (今 / 中間 / 将来)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), Inches(0.4),
             badge, size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left + Inches(0.15), top + Inches(0.4), width - Inches(0.3), Inches(0.4),
             title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # 期間
    add_text(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.3),
             period, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # KPI 行
    add_rect(slide, left + Inches(0.2), top + Inches(1.3),
             width - Inches(0.4), Inches(0.55), LIGHT_GRAY)
    add_text(slide, left + Inches(0.25), top + Inches(1.3),
             (width - Inches(0.5)) / 2, Inches(0.55),
             tenants, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left + (width / 2) + Inches(0.05), top + Inches(1.3),
             (width - Inches(0.5)) / 2, Inches(0.55),
             mrr, size=12, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 本文
    body = '\n'.join(body_lines)
    add_text(slide, left + Inches(0.25), top + Inches(2.0),
             width - Inches(0.5), height - Inches(2.2),
             body, size=11, color=BLACK)


def add_simple_table(slide, left, top, col_widths, row_h, headers, rows,
                     header_color=NAVY, highlight_rows=None):
    """シンプルテーブル描画"""
    if highlight_rows is None:
        highlight_rows = []
    # ヘッダー
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_widths[i], row_h, header_color)
        add_text(slide, x, top, col_widths[i], row_h, h,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_widths[i]
    # 行
    for r, row in enumerate(rows):
        is_hl = r in highlight_rows
        bg = ORANGE if is_hl else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        text_color = WHITE if is_hl else BLACK
        y = top + row_h * (r + 1)
        x = left
        for i, cell in enumerate(row):
            add_rect(slide, x, y, col_widths[i], row_h, bg, line_color=GRAY)
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
            add_text(slide, x + Inches(0.08), y, col_widths[i] - Inches(0.16), row_h,
                     cell, size=11, bold=is_hl, color=text_color,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            x = x + col_widths[i]


# ============================================================
# スライド生成
# ============================================================

def make_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # ============================================================
    # Slide 1: タイトル
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, NAVY)
    add_rect(s, Emu(0), Inches(2.7), SLIDE_W, Inches(0.15), ORANGE)
    add_rect(s, Emu(0), Inches(5.0), SLIDE_W, Inches(0.05), ORANGE)

    add_text(s, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.6),
             'POSLA', size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(0.6),
             '飲食店経済圏ロードマップ',
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.0), Inches(12.5), Inches(0.6),
             '〜 SaaS から、業界インフラへ 〜',
             size=22, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(3.9), Inches(12.5), Inches(0.6),
             '今 (0–6ヶ月)   →   中間 (6ヶ月–2年)   →   将来 (2–5年)',
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.5), Inches(0.5),
             '30店舗・¥500k/月    →    100店舗・¥3M/月    →    1000店舗・¥50M/月',
             size=16, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.4), Inches(12.5), Inches(0.5),
             '対象: 共同経営者・社内チーム / 投資家・銀行',
             size=14, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.5), Inches(0.5),
             '2026年4月版',
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.4),
             'Plus Belief Inc.  |  社外秘 — Confidential',
             size=11, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 2: エグゼクティブサマリー
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'エグゼクティブサマリー', 2)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             'POSLA は飲食店向け SaaS から「飲食店経済圏」のインフラへ進化する',
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    items = [
        '【現在】 オールインワン POS/KDS/在庫/シフト/分析を月1〜5万円で提供する SaaS',
        '【課題】 SaaS 単体では LTV に上限がある (1店舗あたり月額 5万円が天井)',
        '【戦略】 取引データを軸に「フロー全体」を売る経済圏 (発注・採用・決済・送客) を構築',
        '【勝ち筋】 他社 POS は機能を売る → POSLA は飲食店の業務フロー全部を売る',
        '【財務目標】 5年後に ARR ¥6億、テナント 1000店舗、経済圏 GMV ¥50億',
        '【必要資金】 シリーズA で ¥1〜2億 (中間フェーズの営業組織立ち上げに必要)',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.5),
                    items, size=15, line_spacing=1.5)

    add_rect(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.6), NAVY)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.6),
             '"POSLA を入れれば、飲食店経営は POSLA だけで完結する" を5年で実現',
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 3: POSLA の哲学
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'POSLA の根本哲学  〜 機能ではなくフローを売る 〜', 3)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             '他社 POS と POSLA の決定的な違い',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # 左カラム: 他社
    add_rect(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(4.5),
             LIGHT_GRAY, line_color=GRAY)
    add_rect(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.5), GRAY)
    add_text(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.5),
             '他社 POS / KDS の世界',
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    others = [
        'POS は POS だけ',
        'KDS は KDS だけ',
        '在庫管理は別契約',
        'シフトは別 SaaS',
        '予約は別ベンダー',
        '↓',
        'バラバラの月額が積み上がる',
        '画面も操作もデータも分断',
        'スタッフ教育コスト × N',
    ]
    add_bullet_list(s, Inches(0.95), Inches(2.85), Inches(5.4), Inches(3.7),
                    others, size=13, line_spacing=1.3, mark_color=GRAY)

    # 右カラム: POSLA
    add_rect(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(4.5),
             WHITE, line_color=NAVY)
    add_rect(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(0.5), NAVY)
    add_text(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(0.5),
             'POSLA の世界',
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    posla = [
        '予約 → 着席 → 注文 → 厨房',
        '→ 提供 → 会計 → 在庫減算',
        '→ シフト最適化 → 売上分析',
        '→ 顧客台帳 → 次回予約',
        '↓',
        '1つのシステムで業務フロー完結',
        '取引データが全部 POSLA に集まる',
        'AI 提案・経済圏連携の土台に',
    ]
    add_bullet_list(s, Inches(7.1), Inches(2.85), Inches(5.4), Inches(3.7),
                    posla, size=13, line_spacing=1.3, mark_color=ORANGE)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '"フローを売る" = データが集まる = 経済圏化が可能',
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 4: 既存マネタイズの限界
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '既存 SaaS マネタイズの限界', 4)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '月額 SaaS だけでは届かない天井がある',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '月額の上限は 1店舗 ¥5万円。それ以上は飲食店オーナーが受け入れない',
        '機能追加で値上げしても LTV が線形にしか伸びない (×2 が限界)',
        '広告/オウンドメディアの集客だけだと CAC が高止まり (1件 ¥30k〜)',
        '飲食業界の解約率は高い (年間 10%) — SaaS 単体では純増が鈍化',
        '価格競争に巻き込まれる → スマレジ・ぐるなび・USEN 等と同質化',
    ]
    add_bullet_list(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.8),
                    items, size=15, line_spacing=1.4, mark_color=RED)

    add_rect(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.6), LIGHT_GRAY)
    add_text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.4),
             '解決策: 月額 SaaS は "経済圏への入口" として割り切る',
             size=15, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.0),
             '・SaaS 月額: 月¥1〜5万 (固定収益・安定基盤)\n'
             '・経済圏収益: 取引GMVから %手数料 (発注/送客/決済/採用) → 上限なし\n'
             '・"フローを取る"側のビジネスにシフトする',
             size=13, color=BLACK)

    # ============================================================
    # Slide 5: 経済圏戦略の全体像
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '経済圏戦略の全体像', 5)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             'POSLA を中心に「飲食店の取引すべて」を経由させる',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # 中央に POSLA
    cx = Inches(6.667)
    cy = Inches(4.3)
    core_w = Inches(2.4)
    core_h = Inches(1.2)
    add_rect(s, cx - core_w/2, cy - core_h/2, core_w, core_h, NAVY)
    add_text(s, cx - core_w/2, cy - core_h/2, core_w, core_h,
             'POSLA\n取引データ',
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 6方向のサテライト
    satellites = [
        ('飲食店 (テナント)',  Inches(2.0),  Inches(2.0)),
        ('お客さん (来店者)',   Inches(11.3), Inches(2.0)),
        ('発注先 (食材卸)',     Inches(0.5),  Inches(4.3)),
        ('決済 (Stripe)',       Inches(11.8), Inches(4.3)),
        ('採用 (求人)',         Inches(2.0),  Inches(6.0)),
        ('集客 (Google/SNS)',  Inches(11.3), Inches(6.0)),
    ]
    sat_w = Inches(2.4)
    sat_h = Inches(0.8)
    for name, sl, st in satellites:
        add_rect(s, sl, st, sat_w, sat_h, ORANGE, line_color=NAVY)
        add_text(s, sl, st, sat_w, sat_h, name,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '6つのチャネルすべてから手数料を取る = SaaS 単独の 10倍の LTV',
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 6: 3段階ロードマップ概要
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '3段階ロードマップ — 全体像', 6)

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '無理なく積み上げる、5年間の進化シナリオ',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    cards = [
        ('PHASE 1', '今 (基盤固め)', '0–6ヶ月',
         '30店舗', '¥500k / 月',
         ['・SaaS 完成度を上げる',
          '・既存テナントの定着',
          '・予約管理 (L-9) 投入',
          '・LP/申込フロー (L-12)',
          '・最初の10店舗で営業確立',
          '・解約率を10% → 5%に',
          '',
          '▶ ゴール:',
          '  "売れる SaaS" として',
          '  外部に話せる状態']),
        ('PHASE 2', '中間 (データ武器化)', '6ヶ月–2年',
         '100店舗', '¥3M / 月',
         ['・取引データを蓄積',
          '・AI 需要予測の精度向上',
          '・発注先 (食材卸) と提携',
          '・採用プラットフォーム連携',
          '・地域マイクロドミナント',
          '・既存テナント経由で口コミ',
          '',
          '▶ ゴール:',
          '  経済圏の "兆し" を',
          '  投資家に見せられる']),
        ('PHASE 3', '将来 (経済圏完成)', '2–5年',
         '1000店舗', '¥50M / 月',
         ['・POSLA Pay (決済)',
          '・POSLA Order (B2B EC)',
          '・POSLA Jobs (求人)',
          '・POSLA Map (集客)',
          '・全国展開 + 海外進出',
          '・上場 or 事業売却検討',
          '',
          '▶ ゴール:',
          '  "飲食店経済圏" の',
          '  デファクトインフラ']),
    ]

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    card_top = Inches(2.0)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    accent_colors = [BLUE, NAVY, ORANGE]
    for i, (badge, title, period, tenants, mrr, body) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_phase_card(s, left, card_top, card_w, card_h, badge, title, period,
                       tenants, mrr, body, accent_color=accent_colors[i])

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '各フェーズのゴールが達成できないと、次に進まない (リスク管理)',
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 7: PHASE 1 詳細 (今 0-6ヶ月)
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'PHASE 1: 今 (0–6ヶ月)  —  基盤固め', 7,
                  subtitle='30店舗 / ¥500k MRR')

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '"売れる SaaS" を確立し、外部に話せる状態をつくる',
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 左: アクション
    add_text(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.4),
             '【主要アクション】', size=15, bold=True, color=NAVY)
    actions = [
        'プロダクト完成度: P1b 残バグ撲滅 / 予約管理 (L-9) 投入',
        'マーケ基盤: LP + 申込フロー (L-12) / SEO 記事 50本',
        'セルフサーブ申込: Stripe Checkout 24時間自動化',
        '営業 = Hiroが直接10店舗訪問 / フィードバック収集',
        '既存テナント: 月次定例 / NPS計測 / 解約理由ヒアリング',
        '財務: 月次 P/L 整備 / 投資家ピッチ資料 (本資料) 完成',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(4.4),
                    actions, size=12, line_spacing=1.45)

    # 右: KPI
    add_text(s, Inches(7.1), Inches(2.0), Inches(5.6), Inches(0.4),
             '【KPI】', size=15, bold=True, color=NAVY)

    kpi_rows = [
        ['有料テナント数',    '30店舗',     '純増 +5/月'],
        ['MRR',               '¥500k',      '解約 5%以下'],
        ['ARR',               '¥6M',        '12倍まで成長'],
        ['解約率',            '5%/月以下',  '業界平均10%'],
        ['NPS',               '+30以上',    '推奨者比率'],
        ['CAC',               '¥30k以下',   'LTV/CAC ≧ 5'],
    ]
    add_simple_table(s, Inches(7.1), Inches(2.45),
                     [Inches(2.4), Inches(1.5), Inches(1.7)], Inches(0.45),
                     ['指標', '目標', '備考'], kpi_rows)

    add_rect(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5), NAVY)
    add_text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5),
             'GO/NO-GO: 6ヶ月で 30店舗・MRR ¥500k 未達 → PHASE 2 開始を後ろ倒し',
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 8: PHASE 2 詳細 (中間 6ヶ月-2年)
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'PHASE 2: 中間 (6ヶ月–2年)  —  データ武器化', 8,
                  subtitle='100店舗 / ¥3M MRR')

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '取引データを蓄積し、経済圏の "兆し" を投資家に見せる',
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 左: アクション
    add_text(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.4),
             '【主要アクション】', size=15, bold=True, color=NAVY)
    actions = [
        'AI 需要予測: 既存テナント30店舗のデータで精度を実証',
        '発注先連携: 食材卸 1〜3社と提携 (まずは関東地域)',
        '採用連携: バイトル/タウンワーク等と Webhook 連携',
        '営業組織: 営業 2名 + CS 1名 採用 (シリーズA 資金で)',
        '地域ドミナント: 渋谷/新宿/池袋エリアで密度を作る',
        '既存テナント口コミ: 紹介報酬制度 (1店舗紹介 ¥30k)',
        'シリーズA 調達: ¥1〜2億 (PHASE 3 への足場)',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(4.4),
                    actions, size=12, line_spacing=1.4)

    # 右: KPI
    add_text(s, Inches(7.1), Inches(2.0), Inches(5.6), Inches(0.4),
             '【KPI】', size=15, bold=True, color=NAVY)

    kpi_rows = [
        ['有料テナント数',    '100店舗',    '純増 +6/月'],
        ['MRR',               '¥3M',        'SaaS のみ'],
        ['ARR',               '¥36M',       '6倍'],
        ['経済圏 GMV',        '¥1億/年',    '発注+送客'],
        ['経済圏手数料',      '¥3M/年',     '率3%換算'],
        ['解約率',            '3%/月以下',  '改善継続'],
    ]
    add_simple_table(s, Inches(7.1), Inches(2.45),
                     [Inches(2.4), Inches(1.5), Inches(1.7)], Inches(0.45),
                     ['指標', '目標', '備考'], kpi_rows)

    add_rect(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5), NAVY)
    add_text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5),
             'GO/NO-GO: 経済圏売上が ARR の 5%を超えていれば PHASE 3 へ',
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 9: PHASE 3 詳細 (将来 2-5年)
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'PHASE 3: 将来 (2–5年)  —  経済圏完成', 9,
                  subtitle='1000店舗 / ¥50M MRR')

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '飲食店経済圏のデファクトインフラとなり、上場を視野に',
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 左: 経済圏4本柱
    add_text(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.4),
             '【経済圏 4本柱】', size=15, bold=True, color=NAVY)
    actions = [
        'POSLA Pay  — 決済プラットフォーム (Stripe Connect 経由)',
        'POSLA Order — 食材卸 B2B EC マーケットプレイス',
        'POSLA Jobs  — 飲食店特化求人 (シフト×採用統合)',
        'POSLA Map  — お客さん向け予約・送客 (Tabelog 代替)',
        '',
        '※ いずれも既存 POSLA テナントを土台にコールドスタート問題を回避',
        '',
        '【海外展開】 アジア (台湾・タイ・ベトナム) — 2027年〜',
        '【出口戦略】 上場 (マザーズ) or 事業売却 (大手 SaaS / 商社)',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(4.4),
                    actions, size=12, line_spacing=1.4)

    # 右: KPI
    add_text(s, Inches(7.1), Inches(2.0), Inches(5.6), Inches(0.4),
             '【KPI】', size=15, bold=True, color=NAVY)

    kpi_rows = [
        ['有料テナント数',    '1000店舗',  ''],
        ['SaaS MRR',          '¥30M',       '¥360M ARR'],
        ['経済圏 GMV',        '¥50億/年',  ''],
        ['経済圏 手数料収益', '¥250M/年',  '率5%'],
        ['総 ARR',            '¥600M+',    'SaaS+経済圏'],
        ['営業利益率',        '20%以上',    '上場可能水準'],
    ]
    add_simple_table(s, Inches(7.1), Inches(2.45),
                     [Inches(2.4), Inches(1.5), Inches(1.7)], Inches(0.45),
                     ['指標', '目標', '備考'], kpi_rows)

    add_rect(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5), ORANGE)
    add_text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5),
             '5年後: ARR ¥6億 / GMV ¥50億 / 1000店舗 / 営業利益 ¥1.2億',
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 10: 段階別 KPI 一覧
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '段階別 KPI 一覧 (まとめ)', 10)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '5年間で SaaS 単体から経済圏インフラへ',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    headers = ['指標', 'PHASE 1 (6ヶ月)', 'PHASE 2 (2年)', 'PHASE 3 (5年)', '成長倍率']
    rows = [
        ['有料テナント',      '30',        '100',       '1000',      '×33'],
        ['SaaS MRR',         '¥500k',     '¥3M',       '¥30M',      '×60'],
        ['SaaS ARR',         '¥6M',       '¥36M',      '¥360M',     '×60'],
        ['経済圏 GMV',       '—',          '¥1億',      '¥50億',     '×50'],
        ['経済圏手数料',     '—',          '¥3M',       '¥250M',     '×80'],
        ['総 ARR',           '¥6M',       '¥39M',      '¥600M+',    '×100'],
        ['営業利益率',       '— (赤字許容)', '0%',       '20%以上',   ''],
        ['解約率/月',        '5%',        '3%',        '2%',        '改善'],
        ['NPS',              '+30',       '+50',       '+60',       ''],
        ['従業員数',         '3名',       '8名',       '40名',      '×13'],
    ]

    col_w = [Inches(2.4), Inches(2.3), Inches(2.3), Inches(2.3), Inches(1.4)]
    add_simple_table(s, Inches(1.2), Inches(2.2), col_w, Inches(0.42),
                     headers, rows, highlight_rows=[5])

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             '※ PHASE 1 は赤字許容期間 (シード資金で運営)',
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 11: 主要懸念への回答
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '主要懸念への回答', 11)

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             '経営陣 (Hiro) からの2つの本質的な懸念',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Q1
    add_rect(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.5), RED)
    add_text(s, Inches(0.9), Inches(2.0), Inches(11.7), Inches(0.5),
             'Q1.  発注先 (食材卸) の開拓が大変。Hiro が直接行くのか？',
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(2.55), Inches(11.7), Inches(1.7),
             'A. Hiroが直接行く必要はない。3つのレバレッジを使う:\n'
             '   ① 既存テナントの仕入先に POSLA 経由で発注 → 卸側に "POSLA から発注が増える" 実績を見せる\n'
             '   ② 食材卸大手 1社 (例: 国分グループ・三菱食品) と提携 → 1社で数百卸とつながる\n'
             '   ③ 既存テナントから卸を紹介してもらう (紹介料 ¥10k) → 信頼関係付きで提携交渉が短い',
             size=12, color=BLACK, align=PP_ALIGN.LEFT)

    # Q2
    add_rect(s, Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.5), RED)
    add_text(s, Inches(0.9), Inches(4.4), Inches(11.7), Inches(0.5),
             'Q2.  Tabelog 的な集客は鶏卵問題。お客さんが来ないと店が POSLA を辞める',
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(4.95), Inches(11.7), Inches(1.7),
             'A. ゼロから集客サービスを作らない。既存テナントの来店客を初期ユーザーにする:\n'
             '   ① セルフメニュー利用客 (1テナント月平均 800人) × 100店舗 = 月8万人のユーザー基盤\n'
             '   ② Googleレビュー連携 (L-7) で高評価客を Google マップに誘導 → 既存資産を活かす\n'
             '   ③ POSLA Map は PHASE 3 で投入。それまで送客は Google/Tabelog に依存し、自前は作らない',
             size=12, color=BLACK, align=PP_ALIGN.LEFT)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             'いずれも "既存テナント資産" を最大限レバレッジするのが共通戦略',
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 12: 市場規模 (TAM/SAM/SOM)
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '市場規模  (TAM / SAM / SOM)', 12,
                  subtitle='for 投資家')

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '日本の飲食店市場は十分に大きく、シェア 1%でも事業として成立する',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # 3層の同心円風に矩形で表現
    layers = [
        ('TAM (全体)', '日本の飲食店 約65万店', '¥6,500億',
         '市場全体: 全店が POS/業務システムを月10万円使うと仮定', NAVY,
         Inches(0.7), Inches(2.2), Inches(11.9)),
        ('SAM (実質ターゲット)', '導入意欲のある中小チェーン 約10万店', '¥1,000億',
         '実際に POSLA を入れる可能性がある店 — IT リテラシ・規模・業態でフィルタ', BLUE,
         Inches(2.2), Inches(3.5), Inches(8.9)),
        ('SOM (5年で取りに行く)', '1,000店舗 (シェア 1%)', '¥6億 ARR',
         '5年後の現実的な目標 — シェア 1%でも ARR ¥6億の事業', ORANGE,
         Inches(3.7), Inches(4.8), Inches(5.9)),
    ]
    for label, count, size_val, note, color, l, t, w in layers:
        add_rect(s, l, t, w, Inches(1.0), color)
        add_text(s, l + Inches(0.2), t, Inches(2.5), Inches(1.0), label,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, l + Inches(2.7), t, Inches(4.5), Inches(1.0), count,
                 size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, l + w - Inches(2.5), t, Inches(2.3), Inches(1.0), size_val,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.4),
             '出典: 経済産業省 商業動態統計 / 日本フードサービス協会',
             size=11, color=GRAY)
    add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.4),
             '※ TAM 計算: 65万店 × 月¥10万 × 12ヶ月 = ¥7,800億 → 控えめに ¥6,500億で算出',
             size=11, color=GRAY)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '"巨大市場でシェア 1%取る" のが現実的かつ十分大きい',
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 13: 競合比較
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '競合比較', 13, subtitle='for 投資家')

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '"POS だけ" の競合は多いが、"フローを売る" 競合は存在しない',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    headers = ['', 'POSLA', 'スマレジ', 'ぐるなび FineOrder', 'USEN Order']
    rows = [
        ['POS / KDS / 在庫統合',     '◎',  '△',  '○',  '△'],
        ['AI 機能標準装備',           '◎',  '×',  '×',  '×'],
        ['多店舗 / 本部統括',         '◎',  '○',  '○',  '○'],
        ['月額 (1店舗・全機能)',     '¥1〜5万', '¥1〜3万 (機能別)', '¥3〜10万', '¥5〜10万'],
        ['経済圏 (発注/送客)',       '◎ (戦略)', '×',  '○ (送客のみ)', '×'],
        ['データのオープン性',        '◎',  '△',  '×',  '×'],
        ['導入スピード',              '即日 (セルフサーブ)', '1週間', '営業訪問必須', '営業訪問必須'],
    ]

    col_w = [Inches(2.4), Inches(2.3), Inches(2.3), Inches(2.7), Inches(2.5)]
    add_simple_table(s, Inches(0.7), Inches(2.1), col_w, Inches(0.5),
                     headers, rows)

    add_rect(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.45), NAVY)
    add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.45),
             'POSLA の独自性: "オールインワン × AI標準 × 経済圏戦略 × セルフサーブ"',
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 14: ARR 成長予測
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'ARR 成長予測', 14, subtitle='for 投資家')

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             'SaaS + 経済圏のハイブリッド成長',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # 棒グラフ風 (テキストベース)
    years = [
        ('Y0 (現在)',  6,    0,    NAVY),
        ('Y1 (PHASE 1)', 12,  0,    NAVY),
        ('Y2 (PHASE 2)', 36,  3,    BLUE),
        ('Y3',          80,   30,   BLUE),
        ('Y4',          200,  120,  ORANGE),
        ('Y5 (PHASE 3)', 360, 250,  ORANGE),
    ]

    chart_left = Inches(1.0)
    chart_top = Inches(2.2)
    chart_w = Inches(11.3)
    chart_h = Inches(3.8)
    bar_w = chart_w / 7
    max_total = 360 + 250  # Y5 の最大値
    px_per_unit = chart_h.emu / max_total  # 1単位あたりの emu

    # ベースライン
    add_rect(s, chart_left, chart_top + chart_h, chart_w, Emu(20000), BLACK)

    for i, (label, saas, eco, color) in enumerate(years):
        bx = chart_left + bar_w * (i + 0.5) - Inches(0.4)
        # SaaS バー (下)
        saas_h = Emu(int(saas * px_per_unit))
        add_rect(s, bx, chart_top + chart_h - saas_h, Inches(0.8), saas_h, NAVY)
        # 経済圏バー (上)
        if eco > 0:
            eco_h = Emu(int(eco * px_per_unit))
            add_rect(s, bx, chart_top + chart_h - saas_h - eco_h,
                     Inches(0.8), eco_h, ORANGE)
        # ラベル下
        add_text(s, bx - Inches(0.2), chart_top + chart_h + Inches(0.05),
                 Inches(1.2), Inches(0.4), label,
                 size=10, color=BLACK, align=PP_ALIGN.CENTER)
        # 値ラベル上
        total = saas + eco
        add_text(s, bx - Inches(0.2),
                 chart_top + chart_h - Emu(int(total * px_per_unit)) - Inches(0.4),
                 Inches(1.2), Inches(0.35),
                 '¥%dM' % total,
                 size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 凡例
    add_rect(s, Inches(8.5), Inches(6.15), Inches(0.3), Inches(0.25), NAVY)
    add_text(s, Inches(8.85), Inches(6.13), Inches(2.0), Inches(0.3),
             'SaaS ARR', size=11, color=BLACK)
    add_rect(s, Inches(11.0), Inches(6.15), Inches(0.3), Inches(0.25), ORANGE)
    add_text(s, Inches(11.35), Inches(6.13), Inches(2.0), Inches(0.3),
             '経済圏 ARR', size=11, color=BLACK)

    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             'Y3 から経済圏が立ち上がり、Y5 で SaaS と同水準に → ハイブリッドで天井突破',
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 15: ユニットエコノミクス
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'ユニットエコノミクス', 15, subtitle='for 投資家')

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '1テナント当たりの経済性 — 健全な SaaS 指標',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    headers = ['指標', 'PHASE 1', 'PHASE 2', 'PHASE 3', '備考']
    rows = [
        ['ARPU (月額)',          '¥17k',     '¥25k',     '¥30k',     'pro 中心'],
        ['経済圏 ARPU 追加',      '—',        '¥3k',     '¥21k',     '発注/送客手数料'],
        ['総 ARPU',              '¥17k',     '¥28k',     '¥51k',     ''],
        ['解約率 (月)',           '5%',       '3%',       '2%',       '改善前提'],
        ['顧客寿命 (ヶ月)',       '20',       '33',       '50',       '1/解約率'],
        ['LTV',                  '¥340k',    '¥924k',    '¥2.55M',   'ARPU × 寿命'],
        ['CAC',                  '¥30k',     '¥50k',     '¥100k',    '営業組織化で増'],
        ['LTV / CAC',            '11×',      '18×',      '25×',      '健全水準 ≧3×'],
        ['CAC 回収期間',          '2ヶ月',    '2ヶ月',    '2ヶ月',    '優秀'],
        ['粗利率',               '85%',      '80%',      '75%',      'AI/インフラ費用増'],
    ]
    col_w = [Inches(2.4), Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.9)]
    add_simple_table(s, Inches(1.2), Inches(2.0), col_w, Inches(0.42),
                     headers, rows, highlight_rows=[5, 7])

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             'LTV/CAC ≧ 3× が SaaS 健全指標 — POSLA は全フェーズで余裕でクリア',
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 16: リスク分析
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, 'リスク分析と対策', 16)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '主要リスクと、それぞれの緩和策',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    headers = ['リスク', '影響度', '発生確率', '対策']
    rows = [
        ['プロダクト不具合で大規模解約',
         '高', '低', 'P1b 撲滅 / 自動 smoke test / 監視'],
        ['営業組織立ち上げ失敗 (PHASE 2)',
         '高', '中', 'シリーズA 資金で経験者採用 / KPI 短サイクル'],
        ['食材卸との提携が進まない',
         '中', '中', '既存テナント経由 / まず 1 社で実証'],
        ['競合 (スマレジ等) の AI 機能追加',
         '中', '高', '経済圏戦略でロックイン / フロー差別化'],
        ['CAC 高騰 (広告依存)',
         '高', '中', '紹介経由を主軸 / SEO+口コミで非広告'],
        ['資金調達失敗',
         '高', '低', 'PHASE 1 を黒字基調で / 銀行融資並行'],
        ['Hiro 1 人依存 (組織未成熟)',
         '高', '高', 'PHASE 1 でCTO/CSO 採用 / 文書化徹底'],
    ]

    col_w = [Inches(4.0), Inches(1.2), Inches(1.5), Inches(5.6)]
    add_simple_table(s, Inches(0.8), Inches(2.1), col_w, Inches(0.55),
                     headers, rows)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             '最大のリスクは "Hiro 1 人依存"。PHASE 1 で組織化を最優先',
             size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 17: 必要なリソース・体制
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '必要なリソース・体制', 17)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '各フェーズで必要な人員・資金・体制',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    cards = [
        ('PHASE 1', [
            '【人員】 3名',
            '・Hiro (CEO/全般)',
            '・エンジニア 1名',
            '・CS 1名 (バイト可)',
            '',
            '【資金】',
            '自己資金 + シード',
            '¥30M 想定',
            '',
            '【体制】',
            '創業期チーム',
        ]),
        ('PHASE 2', [
            '【人員】 8名',
            '・経営 (Hiro+1)',
            '・エンジニア 3名',
            '・営業 2名',
            '・CS 1名',
            '',
            '【資金】',
            'シリーズA',
            '¥150M 想定',
            '',
            '【体制】',
            '機能組織化',
        ]),
        ('PHASE 3', [
            '【人員】 40名',
            '・経営 5名',
            '・エンジニア 15名',
            '・営業 10名',
            '・CS 5名',
            '・経済圏事業 5名',
            '',
            '【資金】',
            'シリーズB / 銀行',
            '¥500M〜',
            '',
            '【体制】',
            '事業部制',
        ]),
    ]
    card_w = Inches(4.0)
    card_h = Inches(4.5)
    card_top = Inches(2.1)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    accent_colors = [BLUE, NAVY, ORANGE]
    for i, (title, lines) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_rect(s, left + Emu(50000), card_top + Emu(50000), card_w, card_h,
                 RGBColor(0xE0, 0xE0, 0xE0))
        add_rect(s, left, card_top, card_w, card_h, WHITE,
                 line_color=accent_colors[i])
        add_rect(s, left, card_top, card_w, Inches(0.5), accent_colors[i])
        add_text(s, left, card_top, card_w, Inches(0.5), title,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        body = '\n'.join(lines)
        add_text(s, left + Inches(0.25), card_top + Inches(0.65),
                 card_w - Inches(0.5), card_h - Inches(0.8),
                 body, size=12, color=BLACK)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             '※ シリーズA タイミングは PHASE 1 終盤 (12ヶ月時点) を想定',
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 18: 5年後のビジョン
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, NAVY)
    add_rect(s, Emu(0), Inches(2.3), SLIDE_W, Inches(0.1), ORANGE)
    add_footer(s, 18)

    add_text(s, Inches(0.5), Inches(0.6), Inches(12.5), Inches(0.6),
             '5年後のビジョン', size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.5), Inches(0.5),
             '"飲食店経済圏" のデファクトインフラ',
             size=20, color=ORANGE, align=PP_ALIGN.CENTER)

    # 4つの数字
    metrics = [
        ('1,000', '店舗', '導入'),
        ('¥6億', 'ARR', '年商'),
        ('¥50億', 'GMV', '経済圏流通'),
        ('20%', '営業利益率', '上場可能'),
    ]
    met_w = Inches(2.7)
    met_h = Inches(2.0)
    met_top = Inches(2.7)
    met_gap = Inches(0.3)
    met_total = met_w * 4 + met_gap * 3
    met_start = (SLIDE_W - met_total) / 2

    for i, (val, label, sub) in enumerate(metrics):
        left = met_start + (met_w + met_gap) * i
        add_rect(s, left, met_top, met_w, met_h, ORANGE)
        add_text(s, left, met_top + Inches(0.15), met_w, Inches(0.7), val,
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left, met_top + Inches(0.95), met_w, Inches(0.4), label,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left, met_top + Inches(1.4), met_w, Inches(0.4), sub,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.2), Inches(12.5), Inches(0.5),
             '"POSLA を入れれば、飲食店経営は POSLA だけで完結する"',
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.5), Inches(0.5),
             '↓',
             size=24, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.5),
             'これが実現できれば、上場 or 大手による買収',
             size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 19: 次のアクション (90日)
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, WHITE)
    add_title_bar(s, '次の 90日間アクション', 19)

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             'PHASE 1 を始動するための具体的な3ヶ月計画',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    months = [
        ('Month 1', '基盤整備', [
            '・予約管理 (L-9) 設計+実装開始',
            '・LP/申込フロー (L-12) 設計',
            '・本資料を共同経営者・社内に共有',
            '・既存テナント月次定例の仕組み化',
            '・SEO 記事 10本',
        ]),
        ('Month 2', '営業始動', [
            '・予約管理 ベータリリース',
            '・LP 公開 + Stripe Checkout 接続',
            '・Hiro が直接 5店舗営業',
            '・解約理由ヒアリング 30件',
            '・SEO 記事 +20本',
        ]),
        ('Month 3', '量産体制へ', [
            '・予約管理 GA',
            '・新規 5〜10店舗獲得',
            '・投資家向け 1on1 開始',
            '・銀行融資相談',
            '・第二エンジニア面接開始',
        ]),
    ]

    col_w = Inches(4.0)
    col_h = Inches(4.5)
    col_top = Inches(2.0)
    col_gap = Inches(0.25)
    total_w = col_w * 3 + col_gap * 2
    start_left = (SLIDE_W - total_w) / 2

    for i, (badge, title, items) in enumerate(months):
        left = start_left + (col_w + col_gap) * i
        add_rect(s, left + Emu(50000), col_top + Emu(50000), col_w, col_h,
                 RGBColor(0xE0, 0xE0, 0xE0))
        add_rect(s, left, col_top, col_w, col_h, WHITE, line_color=NAVY)
        add_rect(s, left, col_top, col_w, Inches(0.5), ORANGE)
        add_text(s, left, col_top, col_w, Inches(0.5),
                 '%s  —  %s' % (badge, title),
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        body = '\n'.join(items)
        add_text(s, left + Inches(0.25), col_top + Inches(0.7),
                 col_w - Inches(0.5), col_h - Inches(0.9),
                 body, size=12, color=BLACK)

    add_rect(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5), NAVY)
    add_text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5),
             '90日後: 「PHASE 1 を本格スタートできる状態」になっていること',
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 20: クロージング
    # ============================================================
    s = prs.slides.add_slide(blank)
    set_background(s, NAVY)
    add_rect(s, Emu(0), Inches(2.3), SLIDE_W, Inches(0.15), ORANGE)
    add_footer(s, 20)

    add_text(s, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.7),
             'POSLA は SaaS で終わらない',
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(0.5),
             '飲食店経済圏のインフラを目指す',
             size=20, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(3.0), Inches(12.5), Inches(0.6),
             '今:  基盤を固める',
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.7), Inches(12.5), Inches(0.6),
             '中間:  データを武器化する',
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.5), Inches(0.6),
             '将来:  経済圏を完成させる',
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.6), Inches(12.5), Inches(0.5),
             'Discussion / Questions',
             size=18, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.5),
             'Plus Belief Inc.   |   Hiroki Oda',
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

    return prs


def main():
    prs = make_presentation()
    out_dir = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs')
    if not os.path.exists(out_dir):
        os.makedirs(out_dir)
    out_path = os.path.join(out_dir, 'posla-economic-zone-roadmap.pptx')
    prs.save(out_path)
    print('Generated: %s' % out_path)
    print('Slides: %d' % len(prs.slides))


if __name__ == '__main__':
    main()
