#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
POSLA 営業ピッチ プレゼンテーション生成スクリプト
対象: テナント営業 + 関係者(共同経営者・社内チーム・パートナー)
構成: 21スライド、フロー哲学を最初に据える
出力: docs/sales-pitch.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# デザイン定数
# ============================================================
NAVY = RGBColor(0x1A, 0x23, 0x7E)       # ブランドネイビー
ORANGE = RGBColor(0xE6, 0x51, 0x00)     # アクセントオレンジ
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x15, 0x65, 0xC0)

FONT_JP = 'Hiragino Sans'

# 16:9 スライドサイズ (13.333 x 7.5 inch = 1280x720 px)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOTAL_SLIDES = 21

# ============================================================
# ヘルパー関数
# ============================================================

def set_background(slide, color):
    """スライド背景色を設定"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """塗り矩形を追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    """テキストボックス追加"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_footer(slide, page_num, total=TOTAL_SLIDES):
    """フッター: ページ番号 + ブランド"""
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.38), Inches(6), Inches(0.35),
             'POSLA  |  飲食店現場を全部AIで', size=11, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.38), Inches(1.1), Inches(0.35),
             '%d / %d' % (page_num, total), size=11, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_title_bar(slide, title, page_num):
    """コンテンツスライドの上部タイトルバー"""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.0), NAVY)
    add_rect(slide, Emu(0), Inches(1.0), SLIDE_W, Inches(0.08), ORANGE)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             title, size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, page_num)


def add_bullet_list(slide, left, top, width, height, items, size=18, line_spacing=1.4):
    """箇条書きリスト追加"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        # 装飾マーク
        run_mark = p.add_run()
        run_mark.text = '■  '
        run_mark.font.name = FONT_JP
        run_mark.font.size = Pt(size)
        run_mark.font.color.rgb = ORANGE
        run_mark.font.bold = True

        run = p.add_run()
        run.text = item
        run.font.name = FONT_JP
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK
    return box


def add_card(slide, left, top, width, height, title, body_lines, accent_color=NAVY):
    """3カラム用カード"""
    # 影代わりの背面 (薄グレー)
    add_rect(slide, left + Emu(50000), top + Emu(50000), width, height, RGBColor(0xE0, 0xE0, 0xE0))
    # 本体
    add_rect(slide, left, top, width, height, WHITE, line_color=accent_color)
    # アクセントバー
    add_rect(slide, left, top, width, Inches(0.5), accent_color)
    # タイトル
    add_text(slide, left + Inches(0.2), top + Inches(0.05), width - Inches(0.4), Inches(0.4),
             title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 本文
    body = '\n'.join(body_lines)
    add_text(slide, left + Inches(0.25), top + Inches(0.7), width - Inches(0.5), height - Inches(0.9),
             body, size=14, color=BLACK)


def add_flow_box(slide, left, top, width, height, label, color=NAVY):
    """フロー図用の小箱"""
    add_rect(slide, left, top, width, height, color)
    add_text(slide, left, top, width, height, label,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_arrow_text(slide, left, top, width, height, char='→', color=ORANGE, size=20):
    """矢印テキスト"""
    add_text(slide, left, top, width, height, char,
             size=size, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# スライド生成
# ============================================================

def make_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    # ============================================================
    # Slide 1: タイトル
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, NAVY)

    # 装飾オレンジ帯
    add_rect(s, Emu(0), Inches(2.8), SLIDE_W, Inches(0.15), ORANGE)

    add_text(s, Inches(0.5), Inches(0.8), Inches(12.5), Inches(1.0),
             'POSLA', size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.5), Inches(0.7),
             '〜 飲食店現場を全部AIで 〜', size=28, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(3.3), Inches(12.5), Inches(1.2),
             '機能ではなく "業務フロー" を売る\n予約・注文・KDS・会計・在庫・シフト・分析  ぜんぶ込み',
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.1), Inches(12.5), Inches(0.5),
             '1店舗 月20,000円 (税別)  /  全機能込み  /  追加料金なし',
             size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(5.9), Inches(12.5), Inches(0.5),
             '30日間 無料トライアル実施中',
             size=18, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.5),
             'Plus Belief Inc.  |  https://posla.jp',
             size=14, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 2: 飲食店オーナーの3つの悩み
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '飲食店オーナーの「3つの悩み」', 2)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
             'こんなこと、ありませんか？',
             size=22, color=GRAY, align=PP_ALIGN.CENTER)

    cards = [
        ('① 人手不足', [
            '・ホール・キッチン',
            '  両方まわらない',
            '・新人教育に時間がかかる',
            '・スタッフがすぐ辞める',
            '・繁忙期に手が足りない',
        ]),
        ('② システムが分散', [
            '・POS / KDS / 在庫',
            '  / シフト / 予約',
            '  全部バラバラ',
            '・月額が積み上がる',
            '・データ連携できない',
        ]),
        ('③ 経営判断ができない', [
            '・売れ筋がわからない',
            '・原価率が見えない',
            '・どの店が儲かってるか',
            '  分析できない',
            '・データが手元にない',
        ]),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.7)
    card_top = Inches(2.4)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    for i, (title, lines) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_card(s, left, card_top, card_w, card_h, title, lines, accent_color=ORANGE)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
             '原因は "システムを機能ごとに買っているから"。POSLA は違います。',
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 3: POSLA の哲学 — フローを売る
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, 'POSLA の哲学  〜 機能ではなく「フロー」を売る 〜', 3)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             '他社 POS と POSLA の決定的な違い',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # 左カラム: 他社
    add_rect(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(4.5),
             LIGHT_GRAY, line_color=GRAY)
    add_rect(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.5), GRAY)
    add_text(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.5),
             '他社 (機能を売る世界)',
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    others = [
        'POS は POS だけ',
        'KDS は KDS だけ',
        '在庫管理は別契約',
        'シフトは別 SaaS',
        '予約は別ベンダー',
        '↓',
        'バラバラの月額が積み上がる',
        '画面も操作もデータも分断',
        'スタッフ教育コスト × N',
    ]
    add_bullet_list(s, Inches(0.95), Inches(2.85), Inches(5.4), Inches(3.7),
                    others, size=13, line_spacing=1.3)

    # 右カラム: POSLA
    add_rect(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(4.5),
             WHITE, line_color=NAVY)
    add_rect(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(0.5), NAVY)
    add_text(s, Inches(6.85), Inches(2.2), Inches(5.8), Inches(0.5),
             'POSLA (フローを売る世界)',
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    posla_items = [
        '予約 → 着席 → 注文 → 厨房',
        '→ 提供 → 会計 → 在庫減算',
        '→ シフト最適化 → 売上分析',
        '→ 顧客台帳 → 次回予約',
        '↓',
        '1つのシステムで業務フロー完結',
        '取引データが全部 POSLA に集まる',
        'AI 提案・経済圏連携の土台に',
    ]
    add_bullet_list(s, Inches(7.1), Inches(2.85), Inches(5.4), Inches(3.7),
                    posla_items, size=13, line_spacing=1.3)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '"フローを売る" = データが集まる = 経営判断・AI・経済圏化が可能',
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 4: 業務フローの全体像 (ループ図)
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '業務フローの全体像  〜 POSLA がループ全部を担う 〜', 4)

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             'お客さんの来店から次回予約まで、すべて POSLA 1つで完結',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # 上段5箱
    box_w = Inches(2.0)
    box_h = Inches(0.8)
    arrow_w = Inches(0.4)
    row_total = box_w * 5 + arrow_w * 4
    row_left = (SLIDE_W - row_total) / 2

    top_row_y = Inches(2.2)
    upper = ['① 予約', '② 着席', '③ 注文', '④ 厨房', '⑤ 提供']
    for i, label in enumerate(upper):
        x = row_left + (box_w + arrow_w) * i
        add_flow_box(s, x, top_row_y, box_w, box_h, label, color=NAVY)
        if i < 4:
            add_arrow_text(s, x + box_w, top_row_y, arrow_w, box_h, '→')

    # 右側 下向き矢印
    right_arrow_x = row_left + box_w * 5 + arrow_w * 4 - box_w / 2
    add_arrow_text(s, right_arrow_x, top_row_y + box_h, box_w, Inches(0.6), '↓', size=28)

    # 中央バッジ
    badge_w = Inches(3.0)
    badge_h = Inches(1.4)
    badge_left = (SLIDE_W - badge_w) / 2
    badge_top = Inches(3.55)
    add_rect(s, badge_left, badge_top, badge_w, badge_h, ORANGE)
    add_text(s, badge_left, badge_top + Inches(0.2), badge_w, Inches(0.5),
             'POSLA',
             size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, badge_left, badge_top + Inches(0.75), badge_w, Inches(0.5),
             'このフロー全部を1つで',
             size=13, color=WHITE,
             align=PP_ALIGN.CENTER)

    # 左側 上向き矢印
    left_arrow_x = row_left - box_w / 2
    add_arrow_text(s, left_arrow_x, top_row_y + box_h, box_w, Inches(0.6), '↑', size=28)

    # 下段5箱 (右から左へ進む流れ)
    bot_row_y = Inches(5.1)
    lower = ['⑩ 次回予約', '⑨ 顧客台帳', '⑧ 売上分析', '⑦ 在庫減算', '⑥ 会計']
    for i, label in enumerate(lower):
        x = row_left + (box_w + arrow_w) * i
        add_flow_box(s, x, bot_row_y, box_w, box_h, label, color=BLUE)
        if i < 4:
            add_arrow_text(s, x + box_w, bot_row_y, arrow_w, box_h, '←')

    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             '↑ 1つのシステムでこのループを回すから、データが全部つながる',
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
             '"これだけで店が回る" — POSLA の約束',
             size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 5: POSLAとは
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, 'POSLAとは', 5)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
             '飲食店向け オールインワン クラウドPOS  (フローを丸ごと提供)',
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    items = [
        '注文管理 (ハンディPOS) … スタッフがタブレットで注文取り',
        'KDS (キッチンディスプレイ) … 厨房で注文表示・調理進捗管理',
        'POSレジ … 会計・割引・個別会計・現金/カード/QR決済',
        'セルフオーダー … お客さんがQRコードからスマホで注文',
        '在庫・レシピ管理 … 原材料・棚卸し・AI需要予測',
        'シフト管理 … カレンダー・希望提出・AI最適シフト提案',
        '売上レポート … 日報・回転率・スタッフ評価・併売分析',
        'AIアシスタント … SNS生成・売上分析・需要予測・競合調査',
    ]
    add_bullet_list(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.5),
                    items, size=16, line_spacing=1.35)

    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             '↑ ぜんぶ込み。プラン分けなし。追加料金なし。',
             size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 6: 誰向けか
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '誰のためのシステムか', 6)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '機能で分けません。すべてのお客さんに「全機能」をお届けします',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.5), GREEN)
    add_text(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.5),
             '★ 多店舗運営も全機能標準。価格は店舗数で変わります',
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cards = [
        ('個人店オーナー', [
            '・1店舗だけ運営',
            '・人手不足を',
            '  デジタルで補いたい',
            '・売上をデータで',
            '  把握したい',
            '',
            '→ 月額 ¥20,000',
            '   全機能込み',
        ]),
        ('小〜中規模チェーン', [
            '・2〜10店舗運営',
            '・店舗ごとの',
            '  売上比較',
            '・スタッフ統一管理',
            '・本部レポート',
            '',
            '→ 2店舗目以降',
            '   ¥17,000/店舗',
        ]),
        ('大手チェーン', [
            '・10店舗以上',
            '・本部から',
            '  メニューを一括配信',
            '・店舗間ヘルプ',
            '・統合シフト',
            '',
            '→ + ¥3,000/店舗',
            '   本部一括配信',
        ]),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.9)
    card_top = Inches(2.7)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    for i, (title, lines) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_card(s, left, card_top, card_w, card_h, title, lines, accent_color=NAVY)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             '居酒屋・ラーメン・カフェ・焼肉・定食屋 … 業態問わず、1店舗からチェーンまで使えます',
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 7: お客さん側の機能
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '【お客さん側】 セルフオーダー画面', 7)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             'QRコードを読むだけで、自分のスマホから注文できます',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        'QRコード読み取り → メニュー表示 → 注文 → 決済 まで完結',
        '5言語対応 (日本語 / 英語 / 中国語簡体 / 中国語繁体 / 韓国語)',
        'カテゴリ別表示・写真付き・アレルゲン表示・カロリー表示',
        'おすすめ・売り切れ・期間限定の自動表示',
        '【AIウェイター】 自然な会話で「おすすめは？」「辛くない？」に回答',
        '【店員呼び出し】 ボタン1つで店員にチャイム通知',
        '【セルフレジ】 お客さんが自分で支払いも可能',
        '【満足度評価】 食後のアンケート → Googleレビュー誘導',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6),
                    items, size=15, line_spacing=1.35)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '効果: ホールスタッフの注文取り業務が60〜80%削減',
             size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 8: 店員側の機能
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '【店員側】 ホール・キッチン業務', 8)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '注文取り・調理進捗・会計を1つのシステムでスムーズに',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '【ハンディPOS】 タブレットで注文取り → KDSへ即座に転送',
        '【KDS (キッチンディスプレイ)】 注文を厨房タブレットに表示',
        '【AI音声コマンド】 「ラーメン3番できた」と話すだけで完了通知',
        '【ステーション分け】 焼き場・揚げ場・盛り付け 別画面で表示',
        '【テーブル合流・分割】 4人席に2人席を合体、別会計も簡単',
        '【POSレジ】 現金・カード・QR・電子マネー対応',
        '【個別会計】 1テーブルを複数会計に分割',
        '【シフト管理】 GPS打刻・希望提出・AI最適シフト提案',
        '【勤怠】 出退勤・休憩・残業を自動記録',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6),
                    items, size=14, line_spacing=1.3)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '効果: 注文ミス削減・調理時間短縮・新人教育コスト削減',
             size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 9: 経営者側の機能
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '【経営者側】 売上・分析・経営判断', 9)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '数字で見える化、AIで先回り。「勘」ではなく「データ」で経営する',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '【売上レポート】 日報・週報・月報・時間帯別・曜日別',
        '【回転率分析】 テーブル別・時間帯別の回転率',
        '【ABC分析】 売れ筋ランキング・利益貢献度',
        '【スタッフ評価】 売上貢献・接客評価・出勤率',
        '【併売分析】 「ビールと一緒に何が出るか」をAIが分析',
        '【在庫・原材料】 棚卸し・原価率・廃棄率',
        '【AI需要予測】 過去データ + 天気 + 曜日で発注量を提案',
        '【監査ログ】 不正打刻・割引乱用・キャンセル多発を検知',
        '【オーナーダッシュボード】 全店舗の数字を一画面で',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6),
                    items, size=14, line_spacing=1.3)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '効果: 経営判断スピード向上・原価率改善・人件費最適化',
             size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 10: AI機能ハイライト
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, 'POSLAのAI機能 〜 ここが他社と違う 〜', 10)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '飲食店業務に特化したAI機能をすべてのお客さんに標準装備',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '【AIウェイター】 セルフメニュー画面でお客さんの質問に自然な会話で回答',
        '【AI音声コマンド】 厨房でハンズフリー操作 「ラーメン3番できた」',
        '【AI需要予測】 天気・曜日・過去データから明日の発注量を提案',
        '【AI最適シフト】 売上予測 + スタッフ希望から自動でシフト案を生成',
        '【AI SNS生成】 「今日のおすすめ」を入力するとInstagram投稿文を自動生成',
        '【AI売上分析】 「なぜ今週ラーメンが減ったか」をAIが原因分析',
        '【AI競合調査】 周辺店舗のメニュー・価格をGoogle検索から自動収集',
        '【AIヘルプデスク】 マニュアル全文を学習。「シフトの作り方は？」に即答',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6),
                    items, size=14, line_spacing=1.3)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '※ AI APIキーはPOSLAが負担。お客さんの追加コストはゼロ',
             size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 11: 連携機能
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '外部サービス連携', 11)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '既存の決済・POS・レビューサービスとシームレスに連携',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '【Stripe決済】 カード / Apple Pay / Google Pay / Link 標準対応',
        '【Stripeカードリーダー】 BBPOS WisePOS E など実機端末対応',
        '【Stripe Billing】 サブスク自動課金・請求書発行・領収書メール',
        '【スマレジ連携】 既存スマレジから商品マスタを自動インポート',
        '【Googleレビュー連携】 高評価のお客さんを自動でGoogleマップ誘導',
        '【Google Places API】 競合店舗情報の自動取得',
        '【Gemini AI】 自然言語処理・需要予測・分析エンジン',
    ]
    add_bullet_list(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(4.5),
                    items, size=15, line_spacing=1.4)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '※ Stripe の決済手数料はテナント負担 (POSLAの月額には含まれません)',
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 12: プラン構造  〜 1本化 + アドオン 〜
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, 'プラン構造  〜 シンプルに1つだけ 〜', 12)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             'POSLA は機能でプランを分けません。理由は "フローを売っているから"',
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 大きな1プランカード
    plan_left = Inches(1.5)
    plan_top = Inches(2.4)
    plan_w = Inches(10.3)
    plan_h = Inches(2.2)
    add_rect(s, plan_left + Emu(50000), plan_top + Emu(50000), plan_w, plan_h,
             RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(s, plan_left, plan_top, plan_w, plan_h, WHITE, line_color=NAVY)
    add_rect(s, plan_left, plan_top, plan_w, Inches(0.6), NAVY)
    add_text(s, plan_left, plan_top, plan_w, Inches(0.6),
             'POSLA  ─  基本料金 (全機能込み)',
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 価格表示 (横並び)
    add_text(s, plan_left + Inches(0.3), plan_top + Inches(0.85), Inches(4.5), Inches(0.5),
             '1店舗目',
             size=16, bold=True, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_text(s, plan_left + Inches(0.3), plan_top + Inches(1.25), Inches(4.5), Inches(0.7),
             '¥20,000 / 月',
             size=28, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)

    add_text(s, plan_left + Inches(5.3), plan_top + Inches(0.85), Inches(4.7), Inches(0.5),
             '2店舗目以降 (15% OFF)',
             size=16, bold=True, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_text(s, plan_left + Inches(5.3), plan_top + Inches(1.25), Inches(4.7), Inches(0.7),
             '¥17,000 / 店舗',
             size=28, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)

    # オプションカード
    opt_top = Inches(4.85)
    opt_h = Inches(1.4)
    add_rect(s, plan_left + Emu(50000), opt_top + Emu(50000), plan_w, opt_h,
             RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(s, plan_left, opt_top, plan_w, opt_h, WHITE, line_color=ORANGE)
    add_rect(s, plan_left, opt_top, plan_w, Inches(0.5), ORANGE)
    add_text(s, plan_left, opt_top, plan_w, Inches(0.5),
             '+ 本部一括メニュー配信オプション (チェーン向け)',
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, plan_left + Inches(0.3), opt_top + Inches(0.65), plan_w - Inches(0.6), Inches(0.7),
             '+ ¥3,000 / 店舗   (10店舗以上のチェーン本部向け。1店舗運営なら不要)',
             size=18, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '機能で迷わせない・営業で困らせない・契約後も追加請求なし',
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 13: 価格早見表
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '価格早見表 (税別)', 13)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '規模が大きくなるほど 1店舗あたりの単価は下がります',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # 価格表
    table_left = Inches(1.0)
    table_top = Inches(2.1)
    col_widths = [Inches(2.7), Inches(2.5), Inches(3.0), Inches(3.1)]
    row_h = Inches(0.55)

    headers = ['規模', '基本料金', '+本部一括配信', '月額合計']
    rows = [
        ['個人店  1店舗',     '¥20,000',         '—',           '¥20,000'],
        ['小規模  3店舗',     '¥54,000',         '(任意)',      '¥54,000'],
        ['中規模  5店舗',     '¥88,000',         '+¥15,000',    '¥88k〜103k'],
        ['チェーン 10店舗',   '¥173,000',        '+¥30,000',    '¥173k〜203k'],
        ['大手  30店舗',      '¥513,000',        '+¥90,000',    '¥513k〜603k'],
    ]

    # ヘッダー
    x = table_left
    for i, h in enumerate(headers):
        add_rect(s, x, table_top, col_widths[i], row_h, NAVY)
        add_text(s, x, table_top, col_widths[i], row_h, h,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_widths[i]

    # データ行
    for r, row in enumerate(rows):
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        y = table_top + row_h * (r + 1)
        x = table_left
        for i in range(4):
            add_rect(s, x, y, col_widths[i], row_h, bg, line_color=GRAY)
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
            add_text(s, x + Inches(0.1), y, col_widths[i] - Inches(0.2), row_h,
                     row[i], size=13, color=BLACK,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            x = x + col_widths[i]

    # 内訳メモ
    add_text(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.4),
             '【内訳ロジック】',
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.4),
             '・1店舗目  ¥20,000 / 月',
             size=13, color=BLACK)
    add_text(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.4),
             '・2店舗目以降  ¥17,000 / 店舗  (15%引き)',
             size=13, color=BLACK)
    add_text(s, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.4),
             '・本部一括メニュー配信オプション  +¥3,000 / 店舗  (チェーン本部のみ任意)',
             size=13, color=BLACK)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '※ 30日間 全機能無料トライアル付き  /  追加料金・隠れコストなし',
             size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 14: 競合との違い
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '競合と何が違うのか', 14)

    items = [
        '①【フローを売る】 機能の寄せ集めではなく、業務フロー全体を1つで提供',
        '②【オールインワン】 POS + KDS + 在庫 + シフト + 分析が1つに統合',
        '③【AI標準装備】 追加料金なしでAI機能をフル活用 (他社は月+5,000〜)',
        '④【マルチテナント設計】 1店舗から100店舗まで同じシステムで対応',
        '⑤【プラン迷わない】 1プラン+アドオンだけ。営業トーク1秒で完結',
        '⑥【日本の飲食店に特化】 居酒屋・ラーメン・カフェ業態を熟知',
        '⑦【スマレジから移行可】 既存マスタを丸ごとインポート',
        '⑧【現場主導の開発】 実店舗のフィードバックで毎週アップデート',
    ]
    add_bullet_list(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.0),
                    items, size=16, line_spacing=1.4)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             '"フロー全部入りで月¥20k" — これが最大の差別化です',
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 15: 価格比較表
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '他社サービスとの価格比較 (1店舗あたり月額)', 15)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             '同等機能を競合で揃えると…',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    table_left = Inches(1.0)
    table_top = Inches(2.2)
    col_widths = [Inches(4.5), Inches(3.2), Inches(3.6)]
    row_h = Inches(0.55)

    headers = ['サービス', '月額 (1店舗)', '備考']
    rows = [
        ['POSLA  (1店舗目)',         '¥20,000',     'フロー全部入り'],
        ['POSLA  (2店舗目以降)',     '¥17,000',     'ボリューム割引'],
        ['POSLA  本部配信オプション', '+¥3,000',     'チェーン本部のみ任意'],
        ['他社POS A',                '¥15,000〜',  'POSのみ。KDSは別契約'],
        ['他社KDS B',                '¥8,000〜',   'KDSのみ'],
        ['他社在庫管理 C',           '¥12,000〜',  '在庫のみ'],
        ['他社シフト D',             '¥5,000〜',   'シフトのみ'],
        ['↑ 他社を全部足すと',        '¥40,000〜',  '別々の画面・別々の操作', True],
    ]

    # ヘッダー
    x = table_left
    for i, h in enumerate(headers):
        add_rect(s, x, table_top, col_widths[i], row_h, NAVY)
        add_text(s, x, table_top, col_widths[i], row_h, h,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_widths[i]

    # データ行
    for r, row in enumerate(rows):
        is_highlight = (len(row) > 3 and row[3])
        bg = ORANGE if is_highlight else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        text_color = WHITE if is_highlight else BLACK
        y = table_top + row_h * (r + 1)
        x = table_left
        for i in range(3):
            add_rect(s, x, y, col_widths[i], row_h, bg, line_color=GRAY)
            text = row[i]
            bold = is_highlight or (r < 3)
            add_text(s, x + Inches(0.1), y, col_widths[i] - Inches(0.2), row_h,
                     text, size=13, bold=bold, color=text_color,
                     align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x = x + col_widths[i]

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             '"POSLA 1つ" vs "他社4契約" — 単価でも運用負荷でも勝てます',
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 16: 必要な機材
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '導入に必要な機材', 16)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '最小構成なら、タブレット1枚 + プリンター だけで始められます',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    cards = [
        ('最小構成', [
            '・タブレット 1枚',
            '  (iPad / Android可)',
            '・レシートプリンター',
            '  (任意)',
            '',
            '初期費用',
            '5〜10万円',
        ]),
        ('推奨構成', [
            '・タブレット 2〜3枚',
            '  (ホール用 + KDS用)',
            '・レシートプリンター',
            '・QRコード表 (印刷物)',
            '・WiFi環境',
            '初期費用',
            '15〜25万円',
        ]),
        ('フル構成', [
            '・タブレット 5枚以上',
            '・カードリーダー',
            '  (Stripe BBPOS)',
            '・キッチンマイク',
            '  (音声コマンド用)',
            '初期費用',
            '30〜50万円',
        ]),
    ]
    card_w = Inches(4.0)
    card_h = Inches(4.5)
    card_top = Inches(2.2)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_W - total_w) / 2

    for i, (title, lines) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_card(s, left, card_top, card_w, card_h, title, lines, accent_color=NAVY)

    add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             '※ 機材は店舗様にてご準備ください。POSLAから機材販売も可能',
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 17: 導入の流れ
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '導入の流れ', 17)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             'お申込みから運用開始まで、最短3日',
             size=18, color=GRAY, align=PP_ALIGN.CENTER)

    steps = [
        ('Day 1', '無料トライアル申込', 'POSLAサイトで申込\n30日無料・カード登録のみ'),
        ('Day 1', 'メニュー登録', 'スマレジから自動取込\nor 手動入力'),
        ('Day 2', '店舗・テーブル設定', 'フロアマップ作成\nスタッフアカウント発行'),
        ('Day 2', '操作レクチャー', 'オンライン or 訪問\n15分で操作習得'),
        ('Day 3', '実運用開始', 'お客様の注文を取り\n通常営業'),
    ]

    step_w = Inches(2.4)
    step_h = Inches(3.6)
    step_top = Inches(2.3)
    gap = Inches(0.15)
    total_w = step_w * 5 + gap * 4
    start_left = (SLIDE_W - total_w) / 2

    for i, (day, title, body) in enumerate(steps):
        left = start_left + (step_w + gap) * i
        # 番号バッジ
        add_rect(s, left, step_top, step_w, Inches(0.6), ORANGE)
        add_text(s, left, step_top, step_w, Inches(0.6), day,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 本体
        add_rect(s, left, step_top + Inches(0.6), step_w, step_h - Inches(0.6),
                 LIGHT_GRAY, line_color=NAVY)
        add_text(s, left + Inches(0.15), step_top + Inches(0.75), step_w - Inches(0.3), Inches(0.6),
                 title, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.15), step_top + Inches(1.5), step_w - Inches(0.3), Inches(2),
                 body, size=12, color=BLACK, align=PP_ALIGN.CENTER)
        # 矢印 (最後以外)
        if i < 4:
            arrow_left = left + step_w + gap - Inches(0.05)
            add_text(s, arrow_left, step_top + Inches(1.6), Inches(0.3), Inches(0.5),
                     '→', size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             '導入サポート無料・操作レクチャー無料・初期費用なし',
             size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 18: 30日無料トライアル
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '30日間 無料トライアル', 18)

    # 大きな目玉
    add_rect(s, Inches(2.0), Inches(1.4), Inches(9.3), Inches(1.5), ORANGE)
    add_text(s, Inches(2.0), Inches(1.4), Inches(9.3), Inches(1.5),
             '全機能 30日間 完全無料',
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    items = [
        '○ お申込み時に Stripe でクレジットカード登録 (請求は0円)',
        '○ 登録完了で 即日 全機能フル開放 (1プランしかないので迷わない)',
        '○ 30日間 ゆっくり試せる',
        '○ 気に入ったら 何もしなくてOK → 31日目から自動で月額¥20,000課金',
        '○ 解約 or 店舗追加 は マイページから 24時間いつでも',
        '○ 解約後は データ削除 (個人情報も完全削除)',
    ]
    add_bullet_list(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.5),
                    items, size=15, line_spacing=1.45)

    add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             'お申込み: https://posla.jp',
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 19: よくある質問
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, 'よくある質問 TOP3', 19)

    qas = [
        ('Q1. なぜ機能別プランがないのですか？',
         'A. 飲食店の業務は「予約 → 注文 → 会計 → 在庫 → 分析」と全部つながっており、\n'
         '     一部だけ使う運用は本来ありません。機能別に分けると、お客さんが\n'
         '     "何を選べばいいか" で迷ってしまいます。POSLAは "フロー全部" を1価格で提供します。'),
        ('Q2. インターネット回線が不安定でも使えますか？',
         'A. オフライン検知機能があります。回線が切れても注文取りは継続でき、\n'
         '     復旧後に自動同期されます。'),
        ('Q3. スタッフが操作に慣れるまでどのくらい？',
         'A. 平均15分です。注文画面はLINEレベルの直感操作で設計しており、\n'
         '     アルバイトでも初日から使えます。'),
    ]

    y = Inches(1.6)
    for q, a in qas:
        add_text(s, Inches(0.7), y, Inches(11.9), Inches(0.5),
                 q, size=17, bold=True, color=NAVY)
        y = y + Inches(0.5)
        add_text(s, Inches(1.2), y, Inches(11.4), Inches(1.1),
                 a, size=13, color=BLACK)
        y = y + Inches(1.15)

    # ============================================================
    # Slide 20: できないこと (正直開示)
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, WHITE)
    add_title_bar(s, '正直に言うと、現状できないこと', 20)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
             '隠さずお伝えします。今後のロードマップに含まれている機能です',
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    items = [
        '× 予約管理 … 現在開発中。2026年中にリリース予定 (フローの起点)',
        '× セルフレジ (現金) … カード/QRのみ対応。現金セルフレジは未対応',
        '× レシートプリンター直接印刷 … スター/エプソン対応開発中',
        '× LINE公式アカウント連携 … 現状はGoogleレビューのみ対応',
        '× 食べログ予約連携 … API未公開のため未対応',
        '× 軽減税率の自動判定 … 8%/10%は手動設定が必要',
    ]
    add_bullet_list(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.0),
                    items, size=16, line_spacing=1.5)

    add_rect(s, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.6), NAVY)
    add_text(s, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.6),
             '「できる/できない」を最初に正直にお伝えするのがPOSLAの方針です',
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # Slide 21: クロージング
    # ============================================================
    s = prs.slides.add_slide(blank_layout)
    set_background(s, NAVY)
    add_rect(s, Emu(0), Inches(2.5), SLIDE_W, Inches(0.15), ORANGE)
    add_footer(s, 21)

    add_text(s, Inches(0.5), Inches(0.7), Inches(12.5), Inches(1.0),
             'ご検討ありがとうございます',
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(1.8), Inches(12.5), Inches(0.6),
             '機能ではなく、業務フロー全部をお届けします',
             size=22, color=ORANGE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(3.0), Inches(12.5), Inches(0.7),
             'POSLA',
             size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(4.3), Inches(12.5), Inches(0.5),
             '1店舗 月¥20,000  /  全機能込み  /  30日無料',
             size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # コンタクト
    add_text(s, Inches(0.5), Inches(5.3), Inches(12.5), Inches(0.5),
             'お申込み・お問い合わせ',
             size=16, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.5), Inches(0.6),
             'https://posla.jp',
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.4),
             'Plus Belief Inc.',
             size=12, color=ORANGE, align=PP_ALIGN.CENTER)

    return prs


def main():
    prs = make_presentation()
    out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs')
    if not os.path.exists(out_dir):
        os.makedirs(out_dir)
    out_path = os.path.join(out_dir, 'sales-pitch.pptx')
    prs.save(out_path)
    print('Generated: %s' % out_path)
    print('Slides: %d' % len(prs.slides))


if __name__ == '__main__':
    main()
